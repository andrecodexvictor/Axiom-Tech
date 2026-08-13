"""Safe, deterministic extraction of the internal knowledge formats.

Each returned document has normalized text and primitive metadata only, which is
important because Chroma metadata cannot contain nested objects.  PowerPoint is
handled slide-by-slide so citations can point to the original slide.
"""

from __future__ import annotations

import csv
import hashlib
import json
import re
import unicodedata
import zipfile
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional
from xml.etree import ElementTree


class UnsupportedDocumentError(ValueError):
    """Raised for files we intentionally do not ingest."""


class DocumentLoader:
    """Multi-format loader for text, office, structured, web, and PDF documents."""

    SUPPORTED_EXTENSIONS = {
        ".md",
        ".txt",
        ".json",
        ".csv",
        ".xlsx",
        ".pdf",
        ".docx",
        ".html",
        ".htm",
        ".pptx",
    }

    @classmethod
    def load_file(
        cls, file_path: Path, *, source_root: Optional[Path] = None
    ) -> List[Dict[str, Any]]:
        path = Path(file_path)
        if not path.is_file():
            raise FileNotFoundError("File not found: {0}".format(path))
        extension = path.suffix.lower()
        if extension not in cls.SUPPORTED_EXTENSIONS:
            raise UnsupportedDocumentError("Unsupported file type: {0}".format(extension or "(none)"))

        domain = path.parent.name
        if extension in {".md", ".txt"}:
            documents = [cls._document(path, domain, cls._read_text(path), extension)]
        elif extension == ".json":
            documents = [cls._document(path, domain, cls._load_json(path), extension)]
        elif extension == ".csv":
            documents = [cls._document(path, domain, cls._load_csv(path), extension)]
        elif extension == ".xlsx":
            documents = cls._load_excel(path, domain)
        elif extension == ".pdf":
            documents = cls._load_pdf(path, domain)
        elif extension == ".docx":
            documents = [cls._document(path, domain, cls._load_docx(path), extension)]
        elif extension in {".html", ".htm"}:
            documents = [cls._document(path, domain, cls._load_html(path), extension)]
        elif extension == ".pptx":
            documents = cls._load_pptx(path, domain)
        else:  # pragma: no cover - guarded by SUPPORTED_EXTENSIONS
            raise UnsupportedDocumentError("Unsupported file type: {0}".format(extension))

        source_key = cls._source_key(path, source_root)
        for document in documents:
            document["metadata"]["source_key"] = source_key
        return documents

    @classmethod
    def load_directory(cls, directory: Path) -> List[Dict[str, Any]]:
        """Load all supported files recursively, omitting unreadable ones.

        The API service uses ``load_files`` when it needs per-file diagnostics;
        this compatibility method preserves the former CLI interface.
        """

        documents: List[Dict[str, Any]] = []
        for path in cls.iter_supported_files(directory):
            try:
                documents.extend(cls.load_file(path))
            except (OSError, ValueError, zipfile.BadZipFile):
                continue
        return documents

    @classmethod
    def iter_supported_files(cls, directory: Path) -> Iterable[Path]:
        directory = Path(directory)
        if directory.is_file():
            if directory.suffix.lower() in cls.SUPPORTED_EXTENSIONS:
                yield directory
            return
        if not directory.is_dir():
            raise FileNotFoundError("Directory not found: {0}".format(directory))
        for path in sorted(directory.rglob("*")):
            if path.is_file() and not path.name.startswith(".") and path.suffix.lower() in cls.SUPPORTED_EXTENSIONS:
                yield path

    @staticmethod
    def _read_text(path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="ignore")

    @classmethod
    def _document(
        cls,
        path: Path,
        domain: str,
        content: str,
        file_type: str,
        **location: Any
    ) -> Dict[str, Any]:
        metadata: Dict[str, Any] = {
            "source": path.name,
            "source_key": path.as_posix(),
            "domain": domain,
            "file_type": file_type,
            "path": str(path),
        }
        # Keep filesystem identity alongside the content hash.  The source
        # inventory can then detect changed files without reparsing PDFs,
        # spreadsheets, or presentations on every status request.
        try:
            file_stat = path.stat()
            metadata["size_bytes"] = int(file_stat.st_size)
            metadata["modified_ns"] = int(file_stat.st_mtime_ns)
        except OSError:
            # The loader may be used with a transient fixture or virtual path;
            # ingestion remains valid and the inventory will report that the
            # legacy timestamp is unavailable.
            pass
        metadata.update({key: value for key, value in location.items() if value is not None})
        normalized = cls._normalize_text(content)
        metadata["document_hash"] = hashlib.sha256(normalized.encode("utf-8")).hexdigest()
        return {"content": normalized, "metadata": metadata}

    @staticmethod
    def _load_json(path: Path) -> str:
        with path.open("r", encoding="utf-8") as handle:
            return json.dumps(json.load(handle), ensure_ascii=False, indent=2, sort_keys=True)

    @staticmethod
    def _load_csv(path: Path) -> str:
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
            rows = csv.reader(handle)
            return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)

    @classmethod
    def _load_excel(cls, path: Path, domain: str) -> List[Dict[str, Any]]:
        try:
            import openpyxl
        except ImportError as exc:
            raise UnsupportedDocumentError("Excel support requires openpyxl") from exc

        workbook = openpyxl.load_workbook(path, data_only=True, read_only=True)
        try:
            documents: List[Dict[str, Any]] = []
            for worksheet in workbook.worksheets:
                lines = []
                for row in worksheet.iter_rows(values_only=True):
                    values = [
                        str(value).strip()
                        for value in row
                        if value is not None and str(value).strip()
                    ]
                    if values:
                        lines.append(" | ".join(values))
                if lines:
                    documents.append(
                        cls._document(
                            path,
                            domain,
                            "\n".join(lines),
                            path.suffix.lower(),
                            sheet=worksheet.title,
                        )
                    )
            return documents
        finally:
            workbook.close()

    @classmethod
    def _load_pdf(cls, path: Path, domain: str) -> List[Dict[str, Any]]:
        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise UnsupportedDocumentError("PDF support requires pypdf") from exc

        reader = PdfReader(str(path))
        documents: List[Dict[str, Any]] = []
        for page_number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                documents.append(
                    cls._document(
                        path,
                        domain,
                        text,
                        ".pdf",
                        page=page_number,
                        page_count=len(reader.pages),
                    )
                )
        return documents

    @staticmethod
    def _load_docx(path: Path) -> str:
        try:
            import docx
        except ImportError as exc:
            raise UnsupportedDocumentError("DOCX support requires python-docx") from exc
        document = docx.Document(str(path))
        pieces = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    pieces.append(" | ".join(cells))
        return "\n".join(pieces)

    @staticmethod
    def _load_html(path: Path) -> str:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            return path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
        for tag in soup(["script", "style", "noscript", "svg", "iframe"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)

    @classmethod
    def _load_pptx(cls, path: Path, domain: str) -> List[Dict[str, Any]]:
        """Extract a document per slide, with a no-extra-dependency ZIP/XML fallback."""

        slides: List[str] = []
        try:
            from pptx import Presentation

            presentation = Presentation(str(path))
            for slide in presentation.slides:
                text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        text.append(shape.text.strip())
                slides.append("\n".join(text))
        except (ImportError, ValueError, OSError, KeyError):
            slides = cls._load_pptx_xml(path)

        documents: List[Dict[str, Any]] = []
        for slide_number, content in enumerate(slides, start=1):
            if content.strip():
                documents.append(cls._document(path, domain, content, ".pptx", slide=slide_number))
        return documents

    @staticmethod
    def _load_pptx_xml(path: Path) -> List[str]:
        namespace = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
        with zipfile.ZipFile(path) as archive:
            names = sorted(
                (
                    name
                    for name in archive.namelist()
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                ),
                key=DocumentLoader._slide_number,
            )
            slides: List[str] = []
            for name in names:
                root = ElementTree.fromstring(archive.read(name))
                values = [node.text for node in root.iter(namespace) if node.text]
                slides.append("\n".join(values))
            return slides

    @staticmethod
    def _slide_number(name: str) -> int:
        match = re.search(r"slide(\d+)\.xml$", name)
        return int(match.group(1)) if match else 0

    @staticmethod
    def _source_key(path: Path, source_root: Optional[Path]) -> str:
        if source_root is not None:
            try:
                return path.resolve().relative_to(Path(source_root).resolve()).as_posix()
            except (OSError, ValueError):
                pass
        return path.as_posix()

    @staticmethod
    def _normalize_text(content: str) -> str:
        value = unicodedata.normalize("NFC", content).replace("\r\n", "\n").replace("\r", "\n")
        value = "".join(
            character
            for character in value
            if character in {"\n", "\t"} or unicodedata.category(character) != "Cc"
        )
        return re.sub(r"\n{3,}", "\n\n", value).strip()
